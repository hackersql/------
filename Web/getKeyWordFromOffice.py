#coding=utf-8
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import string
import os,sys
import platform
import chardet
import getopt

keyword = ''

def is_text_file(filename):
    text_characters = "".join(map(chr, range(32, 127)) + list(" "))
    _null_trans = string.maketrans("", "")
    blocksize = 512
    
    def istext(s):        
        # Get the non-text characters (maps a character to itself then
        # use the 'remove' option to get rid of the text characters.)
        t = s.translate(_null_trans, text_characters)
    
        # If more than 30% non-text characters, then
        # this is considered a binary file
        if len(s)>0 and float(len(t))/(len(s)) > 0.30:
            return 0
        return 1        
    return istext(open(filename).read(blocksize))
    
def file_encoding(filename):
    with open(filename,'rb') as f:
        data = f.read()
    return chardet.detect(data)["encoding"]

def read_word(filename):
    #打开文档
    document = Document(filename)
    #读取每段资料
    texts = [ paragraph.text for paragraph in document.paragraphs];
    #输出并观察结果，也可以通过其他手段处理文本即可
    flag = 0
    for text in texts:
        if keyword in text:
            if flag == 0:
                print filename,'\n'
                flag = 1
            print text
    #读取表格材料，并输出结果
    tables = [table for table in document.tables];
    for table in tables:
        for row in table.rows:
            for cell in row.cells:
                val = cell.text if isinstance(cell.text,(str,unicode))   else str(cell.text)
                if keyword in val:
                    if flag == 0:
                        print filename,'\n'
                        flag = 1                    
                    print val,'\n',
        
def read_excel(filename):
    wb = load_workbook(filename)
    flag = 0
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.rows:
            for cell in row:
                val = cell.value if isinstance(cell.value,(str,unicode))  else str(cell.value)
                if keyword in val:
                    if flag == 0:
                        print filename,'\n'
                        flag = 1
                    print val,'\n',
                
def read_ppt(filename):
    prs = Presentation(filename)
    flag = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if keyword in run.text:
                        if flag == 0:
                            print filename,'\n'
                            flag = 1
                        print run.text

def read_text(filename):
    encoding = file_encoding(filename)
    if encoding:
        with open(filename,'r') as f:
            line = f.readline()
            flag = 0 
            while line:
                line = line.decode(encoding)
                if keyword in line:
                    if flag == 0:
                        print filename,'\n'
                        flag = 1
                    print line
                line = f.readline()
            


path = ''
keyword = '' 

options, args = getopt.getopt(sys.argv[1:], "hp:k:", ["help", "path=", "keyword="])
for option, value in options:
    if option in ("-h", "--help"):
        print("example: office.py -p . -k 密码")
        sys.exit(0)
    if option in ("-p", "--path"):
        path = value
    if option in ("-k", "--keyword"):
        sysstr = platform.system()
        if(sysstr =="Windows"):
            keyword = value.decode('gbk')
        else:
            keyword = value.decode('utf8')
        
if not path or not keyword:
    print u'参数不正确\n'
    print("example: office.py -p . -k 密码")
    sys.exit(0)
    
files = list()
for dirpath,dirnames,filenames in os.walk(path):
    for filename in filenames:
        fullpath=os.path.join(dirpath,filename)
        files.append(fullpath)
        
for filename in files:
    extend = os.path.splitext(filename)[1]
    try:
        if extend == '.docx':
            read_word(filename)
        if 'xls' in extend:
            read_excel(filename)
        if extend == '.pptx':
            read_ppt(filename)
        if is_text_file(filename):
            read_text(filename)
    except Exception,e:  
        print Exception,":",e
        